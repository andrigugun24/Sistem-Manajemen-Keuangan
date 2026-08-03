"""
Generate PowerPoint Presentation - Minimalist Business Style
Based on: Minimalist Business Slides.pptx template
Skripsi: Andri Gugun - Sistem Informasi Manajemen Keuangan Terintegrasi
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy, os

# ============================================================
# Create fresh presentation and import theme from template
# ============================================================
TEMPLATE = r"C:\wamp64\www\FinalProject\Minimalist Business Slides.pptx"

# Load the template to grab its slide layouts
template_prs = Presentation(TEMPLATE)

# Create a new blank presentation with the same dimensions
prs = Presentation()
prs.slide_width = template_prs.slide_width
prs.slide_height = template_prs.slide_height

# ============================================================
# COLOR PALETTE (from template analysis)
# ============================================================
CLR_DARK      = RGBColor(0x0B, 0x06, 0x48)   # Dark navy from template
CLR_BLUE_MAIN = RGBColor(0x5F, 0x7D, 0x95)   # Main blue-gray
CLR_BLUE_DARK = RGBColor(0x44, 0x5D, 0x73)   # Darker blue
CLR_BLUE_MED  = RGBColor(0x79, 0x94, 0xA9)   # Medium blue
CLR_BLUE_LT   = RGBColor(0x86, 0x9F, 0xB1)   # Light blue
CLR_SILVER    = RGBColor(0xA5, 0xB7, 0xC6)   # Silver
CLR_CLOUD     = RGBColor(0xC9, 0xD4, 0xDC)   # Cloud gray
CLR_SNOW      = RGBColor(0xE3, 0xE9, 0xED)   # Snow
CLR_CREAM     = RGBColor(0xF5, 0xF2, 0xEE)   # Cream/warm white bg
CLR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
CLR_BLACK     = RGBColor(0x20, 0x21, 0x24)   # Near black for text
CLR_GRAY      = RGBColor(0x4D, 0x4D, 0x4D)   # Body text gray
CLR_LIGHT_GRAY= RGBColor(0x9E, 0x9E, 0x9E)   # Muted text

# Fonts from template
FONT_TITLE = "Vidaloka"
FONT_BODY  = "Montserrat"

TOTAL_SLIDES = 13

# ============================================================
# HELPER FUNCTIONS
# ============================================================

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False,
                color=CLR_BLACK, alignment=PP_ALIGN.LEFT, font_name=FONT_BODY,
                italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_textbox(slide, left, top, width, height, runs_data, alignment=PP_ALIGN.LEFT):
    """Add textbox with multiple styled runs. runs_data = [(text, size, bold, color, font_name), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    for i, (text, size, bold, color, font_name) in enumerate(runs_data):
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txBox

def add_multiline(slide, left, top, width, height, lines_data, alignment=PP_ALIGN.LEFT, line_spacing=Pt(6)):
    """lines_data = [(text, font_size, bold, color, font_name), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color, font_name) in enumerate(lines_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        p.space_after = line_spacing
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font_name
    return txBox

def add_shape_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_slide_number(slide, num):
    add_textbox(slide, Inches(8.7), Inches(4.85), Inches(1.0), Inches(0.3),
                f"{num}/{TOTAL_SLIDES}", font_size=8, color=CLR_LIGHT_GRAY,
                alignment=PP_ALIGN.RIGHT)

def add_page_title(slide, title, subtitle=None):
    """Standard page title at top-left like the template."""
    add_textbox(slide, Inches(0.78), Inches(0.45), Inches(8.44), Inches(0.63),
                title, font_size=22, bold=False, color=CLR_BLACK,
                font_name=FONT_TITLE, alignment=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, Inches(0.78), Inches(1.05), Inches(8.44), Inches(0.35),
                    subtitle, font_size=11, color=CLR_BLUE_MAIN,
                    font_name=FONT_BODY, italic=True)

def add_bullet_items(slide, left, top, width, height, items, font_size=11,
                     color=CLR_GRAY, bullet_char="\u2022"):
    """Minimalist bullet list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = FONT_BODY
        p.space_after = Pt(5)
    
    return txBox

def add_number_badge(slide, left, top, number, color=CLR_BLUE_MAIN, size=Inches(0.4)):
    """Small circular number badge."""
    circle = add_circle(slide, left, top, size, color)
    add_textbox(slide, left, top + Inches(0.04), size, Inches(0.35),
                str(number), font_size=14, bold=True, color=CLR_WHITE,
                alignment=PP_ALIGN.CENTER, font_name=FONT_BODY)

# Use BLANK layout from the new presentation
blank_layout = prs.slide_layouts[6]  # Standard blank layout


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, CLR_CREAM)

# Decorative left bar
add_shape_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(5.63), CLR_BLUE_MAIN)

# Top accent line
add_shape_rect(slide, Inches(0.78), Inches(1.0), Inches(1.5), Pt(2), CLR_BLUE_MAIN)

# "SIDANG SKRIPSI" label
add_textbox(slide, Inches(0.78), Inches(1.15), Inches(3.0), Inches(0.35),
            "SIDANG SKRIPSI", font_size=11, bold=True, color=CLR_BLUE_MAIN,
            font_name=FONT_BODY)

# Main Title
add_multiline(slide, Inches(0.78), Inches(1.55), Inches(7.5), Inches(2.2), [
    ("Rancang Bangun Sistem Informasi", 24, False, CLR_BLACK, FONT_TITLE),
    ("Manajemen Keuangan Terintegrasi", 24, False, CLR_BLACK, FONT_TITLE),
    ("Berbasis Web pada Pondok Pesantren", 24, False, CLR_BLACK, FONT_TITLE),
    ("Latahzan Citeras", 24, False, CLR_BLACK, FONT_TITLE),
], alignment=PP_ALIGN.LEFT, line_spacing=Pt(4))

# Subtitle tech
add_textbox(slide, Inches(0.78), Inches(3.35), Inches(7.5), Inches(0.35),
            "Menggunakan Framework Laravel dan React", font_size=13,
            color=CLR_BLUE_MAIN, font_name=FONT_BODY, italic=True)

# Bottom line
add_shape_rect(slide, Inches(0.78), Inches(3.85), Inches(3.0), Pt(1), CLR_CLOUD)

# Author info
add_multiline(slide, Inches(0.78), Inches(4.0), Inches(4.5), Inches(1.3), [
    ("Disusun Oleh:", 9, False, CLR_LIGHT_GRAY, FONT_BODY),
    ("Andri Gugun", 16, False, CLR_BLACK, FONT_TITLE),
    ("NIM: [Isi NIM Anda]", 10, False, CLR_GRAY, FONT_BODY),
    ("Dosen Pembimbing: [Isi Nama Dosen]", 10, False, CLR_GRAY, FONT_BODY),
], alignment=PP_ALIGN.LEFT, line_spacing=Pt(3))

# Right side decorative circle
add_circle(slide, Inches(7.5), Inches(0.6), Inches(2.5), CLR_SNOW)
add_circle(slide, Inches(7.8), Inches(0.9), Inches(1.9), CLR_CLOUD)


# ============================================================
# SLIDE 2: LATAR BELAKANG
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, CLR_WHITE)

add_page_title(slide, "Latar Belakang")

# Three content cards
cards = [
    ("Kondisi Saat Ini", CLR_BLUE_MAIN, [
        "Pondok Pesantren Latahzan Citeras memiliki aktivitas transaksi yang beragam",
        "Meliputi SPP, tabungan santri, gaji guru, dan arus kas operasional",
        "Pencatatan masih mengandalkan buku besar fisik secara manual",
    ]),
    ("Masalah yang Timbul", CLR_BLUE_DARK, [
        "Rentan terjadi human error dalam perhitungan manual",
        "Proses rekapitulasi lambat, laporan sering terlambat",
        "Orang tua harus menunggu lama untuk validasi riwayat pembayaran",
    ]),
    ("Urgensi Digitalisasi", CLR_BLUE_MED, [
        "Keterlambatan birokrasi menghambat pengambilan kebijakan taktis pimpinan yayasan",
        "Fragmentasi data menyulitkan audit dan transparansi keuangan",
        "Migrasi ke ekosistem digital menjadi sangat mendesak",
    ]),
]

for i, (title, color, items) in enumerate(cards):
    x = Inches(0.6) + Inches(i * 3.0)
    y = Inches(1.45)
    w = Inches(2.75)
    
    # Card accent line on top
    add_shape_rect(slide, x, y, w, Pt(3), color)
    
    # Card title
    add_textbox(slide, x + Inches(0.05), y + Inches(0.15), w - Inches(0.1), Inches(0.4),
                title, font_size=13, bold=True, color=color, font_name=FONT_TITLE)
    
    # Card items
    add_bullet_items(slide, x + Inches(0.05), y + Inches(0.55), w - Inches(0.1), Inches(2.8),
                     items, font_size=9, color=CLR_GRAY)

add_slide_number(slide, 2)


# ============================================================
# SLIDE 3: RUMUSAN MASALAH & TUJUAN
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, CLR_WHITE)

add_page_title(slide, "Rumusan Masalah & Tujuan Penelitian")

# Left column - Rumusan Masalah
add_shape_rect(slide, Inches(0.78), Inches(1.4), Inches(4.0), Pt(3), CLR_BLUE_MAIN)
add_textbox(slide, Inches(0.78), Inches(1.55), Inches(4.0), Inches(0.4),
            "Rumusan Masalah", font_size=15, bold=False, color=CLR_BLUE_MAIN,
            font_name=FONT_TITLE)

rumusan = [
    "Bagaimana memformulasikan arsitektur sistem digital untuk memangkas kelambatan birokrasi pembukuan konvensional?",
    "Bagaimana menerapkan komputasi otomatis untuk menekan human error?",
    "Bagaimana memfasilitasi transparansi informasi finansial bagi wali santri?",
]

for i, item in enumerate(rumusan):
    y = Inches(2.05) + Inches(i * 0.85)
    add_number_badge(slide, Inches(0.78), y, i+1, CLR_BLUE_MAIN)
    add_textbox(slide, Inches(1.3), y, Inches(3.4), Inches(0.75),
                item, font_size=9, color=CLR_GRAY)

# Right column - Tujuan
add_shape_rect(slide, Inches(5.2), Inches(1.4), Inches(4.0), Pt(3), CLR_BLUE_DARK)
add_textbox(slide, Inches(5.2), Inches(1.55), Inches(4.0), Inches(0.4),
            "Tujuan Penelitian", font_size=15, bold=False, color=CLR_BLUE_DARK,
            font_name=FONT_TITLE)

tujuan = [
    "Mentransformasi efisiensi manajemen pembukuan dari manual ke digital",
    "Mengimplementasikan algoritma otomatis untuk akurasi perhitungan saldo kas",
    "Menyediakan media transparansi yang valid bagi orang tua santri dan pimpinan yayasan",
]

for i, item in enumerate(tujuan):
    y = Inches(2.05) + Inches(i * 0.85)
    add_number_badge(slide, Inches(5.2), y, i+1, CLR_BLUE_DARK)
    add_textbox(slide, Inches(5.72), y, Inches(3.4), Inches(0.75),
                item, font_size=9, color=CLR_GRAY)

add_slide_number(slide, 3)


# ============================================================
# SLIDE 4: BATASAN MASALAH
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, CLR_WHITE)

add_page_title(slide, "Batasan Masalah", "Ruang lingkup dan fokus penelitian")

batasan = [
    ("Ruang Lingkup", "Dioperasikan secara eksklusif untuk kebutuhan internal Pondok Pesantren Latahzan Citeras.", CLR_BLUE_MAIN),
    ("Fitur Inti", "Manajemen SPP bulanan, tabungan santri, penggajian tenaga pendidik, dan laporan pos kas operasional.", CLR_BLUE_DARK),
    ("Teknologi", "Arsitektur Web (SPA) dengan Laravel (Backend), React + Inertia.js (Frontend), dan MySQL (Database).", CLR_BLUE_MED),
    ("Pembatasan", "Belum menggunakan payment gateway pihak ketiga. Verifikasi bukti transfer masih dilakukan manual oleh bendahara.", CLR_BLUE_LT),
]

for i, (title, desc, color) in enumerate(batasan):
    col = i % 2
    row = i // 2
    x = Inches(0.78) + Inches(col * 4.5)
    y = Inches(1.65) + Inches(row * 1.6)
    
    # Accent bar
    add_shape_rect(slide, x, y, Pt(4), Inches(1.2), color)
    
    # Title
    add_textbox(slide, x + Inches(0.2), y, Inches(3.8), Inches(0.35),
                title, font_size=14, bold=False, color=color, font_name=FONT_TITLE)
    
    # Description
    add_textbox(slide, x + Inches(0.2), y + Inches(0.4), Inches(3.8), Inches(0.8),
                desc, font_size=9.5, color=CLR_GRAY)

add_slide_number(slide, 4)


# ============================================================
# SLIDE 5: METODOLOGI PENELITIAN
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, CLR_WHITE)

add_page_title(slide, "Metodologi Penelitian")

# Pengumpulan Data
add_shape_rect(slide, Inches(0.78), Inches(1.4), Inches(3.5), Pt(3), CLR_BLUE_MAIN)
add_textbox(slide, Inches(0.78), Inches(1.55), Inches(3.5), Inches(0.35),
            "Pengumpulan Data", font_size=14, bold=False, color=CLR_BLUE_MAIN, font_name=FONT_TITLE)

add_bullet_items(slide, Inches(0.78), Inches(1.95), Inches(3.5), Inches(1.2), [
    "Observasi: Workflow pencatatan manual di pesantren",
    "Wawancara: Dengan bendahara dan pengurus yayasan",
    "Studi Pustaka: Referensi jurnal dan buku terkait",
], font_size=9.5, color=CLR_GRAY)

# Metode Pengembangan
add_shape_rect(slide, Inches(5.2), Inches(1.4), Inches(4.0), Pt(3), CLR_BLUE_DARK)
add_textbox(slide, Inches(5.2), Inches(1.55), Inches(4.0), Inches(0.35),
            "Metode Pengembangan: Agile", font_size=14, bold=False, color=CLR_BLUE_DARK, font_name=FONT_TITLE)

add_textbox(slide, Inches(5.2), Inches(1.95), Inches(4.0), Inches(0.6),
            "Menggunakan Agile Development yang adaptif terhadap perubahan kebutuhan secara iterasi bertahap.",
            font_size=9.5, color=CLR_GRAY)

# Agile Phases - horizontal flow
phases = [
    ("Plan", "Formulasi\nKebutuhan"),
    ("Design", "Pemodelan\nUML & DB"),
    ("Develop", "Konstruksi\nKode"),
    ("Test", "Validasi\nMutu"),
    ("Deploy", "Rilis\nSistem"),
    ("Review", "Evaluasi\n& Iterasi"),
]

# Background band
add_shape_rect(slide, Inches(0.5), Inches(3.4), Inches(9.0), Inches(1.6), CLR_SNOW)

for i, (name, desc) in enumerate(phases):
    x = Inches(0.65) + Inches(i * 1.5)
    
    # Circle with number
    add_circle(slide, x + Inches(0.25), Inches(3.5), Inches(0.45), CLR_BLUE_MAIN)
    add_textbox(slide, x + Inches(0.25), Inches(3.55), Inches(0.45), Inches(0.35),
                str(i+1), font_size=14, bold=True, color=CLR_WHITE,
                alignment=PP_ALIGN.CENTER, font_name=FONT_BODY)
    
    # Phase name
    add_textbox(slide, x, Inches(4.05), Inches(0.95), Inches(0.3),
                name, font_size=11, bold=True, color=CLR_BLUE_DARK,
                alignment=PP_ALIGN.CENTER, font_name=FONT_BODY)
    
    # Phase desc
    add_textbox(slide, x - Inches(0.1), Inches(4.3), Inches(1.15), Inches(0.6),
                desc, font_size=8, color=CLR_LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER)
    
    # Arrow connector
    if i < len(phases) - 1:
        add_textbox(slide, x + Inches(0.95), Inches(3.55), Inches(0.5), Inches(0.35),
                    "\u2192", font_size=14, color=CLR_CLOUD,
                    alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 5)


# ============================================================
# SLIDE 6: ANALISIS SISTEM BERJALAN vs USULAN
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, CLR_WHITE)

add_page_title(slide, "Sistem Berjalan vs Sistem Usulan")

# Left - Sistem Berjalan
add_rounded_rect(slide, Inches(0.6), Inches(1.35), Inches(4.2), Inches(3.5), CLR_SNOW, CLR_CLOUD)
add_textbox(slide, Inches(0.85), Inches(1.45), Inches(3.7), Inches(0.45),
            "Sistem Berjalan (Manual)", font_size=14, bold=False, color=CLR_BLUE_DARK, font_name=FONT_TITLE)

# Red x marks
berjalan_items = [
    "Pencatatan SPP, tabungan, dan gaji dilakukan secara manual (kalkulator & kertas)",
    "Arsip buku besar rentan hilang atau rusak",
    "Proses rekapitulasi lambat dan rentan human error",
    "Tidak ada transparansi informasi bagi orang tua santri",
]
for i, item in enumerate(berjalan_items):
    y = Inches(2.0) + Inches(i * 0.65)
    add_rich_textbox(slide, Inches(0.85), y, Inches(3.7), Inches(0.55), [
        ("\u2717  ", 11, True, RGBColor(0xC0, 0x39, 0x2B), FONT_BODY),
        (item, 9, False, CLR_GRAY, FONT_BODY),
    ])

# Right - Sistem Usulan
add_rounded_rect(slide, Inches(5.2), Inches(1.35), Inches(4.2), Inches(3.5), CLR_SNOW, CLR_CLOUD)
add_textbox(slide, Inches(5.45), Inches(1.45), Inches(3.7), Inches(0.45),
            "Sistem Usulan (Digital)", font_size=14, bold=False, color=CLR_BLUE_MAIN, font_name=FONT_TITLE)

# Green check marks
usulan_items = [
    "Sentralisasi data dalam satu database MySQL yang aman",
    "Otomatisasi perhitungan tagihan SPP, saldo tabungan, dan slip gaji",
    "Generate dokumen nota/kuitansi dan slip gaji secara digital (PDF)",
    "Sistem hak akses Multi-user terintegrasi",
]
for i, item in enumerate(usulan_items):
    y = Inches(2.0) + Inches(i * 0.65)
    add_rich_textbox(slide, Inches(5.45), y, Inches(3.7), Inches(0.55), [
        ("\u2713  ", 11, True, RGBColor(0x27, 0xAE, 0x60), FONT_BODY),
        (item, 9, False, CLR_GRAY, FONT_BODY),
    ])

add_slide_number(slide, 6)


# ============================================================
# SLIDE 7: PERANCANGAN SISTEM (Use Case)
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, CLR_WHITE)

add_page_title(slide, "Perancangan Sistem", "Use Case Diagram & Aktor Utama")

# Placeholder for Use Case Diagram
add_rounded_rect(slide, Inches(0.78), Inches(1.45), Inches(4.5), Inches(3.3), CLR_SNOW, CLR_CLOUD)
add_textbox(slide, Inches(1.5), Inches(2.7), Inches(3.0), Inches(0.5),
            "[ Sisipkan Use Case Diagram ]", font_size=11, color=CLR_LIGHT_GRAY,
            alignment=PP_ALIGN.CENTER, italic=True)

# Actor cards on the right
actors = [
    ("Admin", "Kelola data master (siswa, kelas, guru), akun pengguna, dan pengaturan sistem", CLR_BLUE_MAIN),
    ("Bendahara", "Mengelola tagihan, proses pembayaran, mutasi tabungan, transaksi kas dan gaji", CLR_BLUE_DARK),
    ("Orang Tua", "Melihat rincian tagihan dan riwayat pembayaran anak secara real-time", CLR_BLUE_MED),
    ("Kepsek & Yayasan", "Memantau laporan dan statistik keuangan via Dashboard Analitik", CLR_BLUE_LT),
]

for i, (name, desc, color) in enumerate(actors):
    y = Inches(1.45) + Inches(i * 0.82)
    
    # Accent bar
    add_shape_rect(slide, Inches(5.6), y, Pt(3), Inches(0.7), color)
    
    # Name
    add_textbox(slide, Inches(5.8), y, Inches(3.5), Inches(0.3),
                name, font_size=11, bold=True, color=color, font_name=FONT_BODY)
    
    # Description
    add_textbox(slide, Inches(5.8), y + Inches(0.28), Inches(3.5), Inches(0.42),
                desc, font_size=8, color=CLR_GRAY)

add_slide_number(slide, 7)


# ============================================================
# SLIDE 8: PERANCANGAN BASIS DATA
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, CLR_WHITE)

add_page_title(slide, "Perancangan Basis Data", "Struktur database dan relasi antar entitas")

# Placeholder for ERD
add_rounded_rect(slide, Inches(0.78), Inches(1.45), Inches(5.5), Inches(2.0), CLR_SNOW, CLR_CLOUD)
add_textbox(slide, Inches(2.2), Inches(2.2), Inches(3.0), Inches(0.4),
            "[ Sisipkan ERD / LRS di sini ]", font_size=11, color=CLR_LIGHT_GRAY,
            alignment=PP_ALIGN.CENTER, italic=True)

# 4 Module segments on the right
segments = [
    ("Autentikasi", "Users, roles,\npermissions", CLR_BLUE_MAIN),
    ("Akademik", "Siswa, kelas,\nguru, jabatan", CLR_BLUE_DARK),
    ("Keuangan", "Tagihan, jenis,\npembayaran", CLR_BLUE_MED),
    ("Kas & Tabungan", "Transaksi kas,\ntabungan, pos", CLR_BLUE_LT),
]

for i, (title, desc, color) in enumerate(segments):
    y = Inches(1.45) + Inches(i * 0.5)
    x = Inches(6.6)
    
    add_shape_rect(slide, x, y, Pt(3), Inches(0.4), color)
    add_textbox(slide, x + Inches(0.12), y, Inches(1.2), Inches(0.25),
                title, font_size=9, bold=True, color=color, font_name=FONT_BODY)
    add_textbox(slide, x + Inches(1.35), y, Inches(1.5), Inches(0.4),
                desc.replace('\n', ', '), font_size=8, color=CLR_GRAY)

# Highlight: Polimorfik relation
add_shape_rect(slide, Inches(0.78), Inches(3.7), Inches(8.44), Pt(2), CLR_BLUE_MAIN)
add_textbox(slide, Inches(0.78), Inches(3.85), Inches(8.44), Inches(0.35),
            "Keunggulan Desain: Relasi Polimorfik", font_size=13, bold=False,
            color=CLR_BLUE_MAIN, font_name=FONT_TITLE)
add_textbox(slide, Inches(0.78), Inches(4.2), Inches(8.44), Inches(0.6),
            "Tabel transaksi_kas menggunakan pendekatan Relasi Polimorfik agar dapat mencatat arus kas secara dinamis dari berbagai sumber transaksi (pembayaran SPP, pengeluaran gaji, dsb.) tanpa redundansi data.",
            font_size=9.5, color=CLR_GRAY)

add_slide_number(slide, 8)


# ============================================================
# SLIDE 9: IMPLEMENTASI UI
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, CLR_WHITE)

add_page_title(slide, "Implementasi Antarmuka", "Single Page Application \u2014 React + Inertia.js")

# 3 Screenshot placeholders
screenshots = [
    ("Dashboard Analitik", "Grafik tren pemasukan\nvs pengeluaran dan\nstatus kepatuhan SPP"),
    ("Form Pembayaran", "Proses pencatatan\npembayaran SPP dengan\nvalidasi & bukti transfer"),
    ("Cetak Laporan (PDF)", "Export dan cetak slip\ngaji, bukti kuitansi,\ndan buku kas umum"),
]

for i, (title, desc) in enumerate(screenshots):
    x = Inches(0.6) + Inches(i * 3.1)
    
    # Screenshot placeholder
    add_rounded_rect(slide, x, Inches(1.45), Inches(2.8), Inches(2.0), CLR_SNOW, CLR_CLOUD)
    add_textbox(slide, x + Inches(0.3), Inches(2.2), Inches(2.2), Inches(0.4),
                "[ Screenshot ]", font_size=10, color=CLR_LIGHT_GRAY,
                alignment=PP_ALIGN.CENTER, italic=True)
    
    # Title
    add_textbox(slide, x + Inches(0.1), Inches(3.55), Inches(2.6), Inches(0.3),
                title, font_size=11, bold=True, color=CLR_BLUE_MAIN,
                alignment=PP_ALIGN.CENTER, font_name=FONT_BODY)
    
    # Description
    add_textbox(slide, x + Inches(0.1), Inches(3.85), Inches(2.6), Inches(0.7),
                desc.replace('\n', ' '), font_size=8.5, color=CLR_GRAY,
                alignment=PP_ALIGN.CENTER)

# Bottom note
add_textbox(slide, Inches(0.78), Inches(4.65), Inches(8.44), Inches(0.3),
            "Dibangun dengan pendekatan SPA yang responsif dan cepat \u2014 siap untuk Live Demo",
            font_size=9, color=CLR_BLUE_MAIN, alignment=PP_ALIGN.CENTER, italic=True)

add_slide_number(slide, 9)


# ============================================================
# SLIDE 10: PENGUJIAN SISTEM
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, CLR_WHITE)

add_page_title(slide, "Pengujian Sistem", "Black-Box Testing & White-Box Testing")

# Black-Box
add_shape_rect(slide, Inches(0.78), Inches(1.5), Inches(4.0), Pt(3), CLR_BLUE_MAIN)
add_textbox(slide, Inches(0.78), Inches(1.65), Inches(4.0), Inches(0.35),
            "Black-Box Testing", font_size=14, bold=False, color=CLR_BLUE_MAIN, font_name=FONT_TITLE)

add_bullet_items(slide, Inches(0.78), Inches(2.05), Inches(4.0), Inches(1.5), [
    "Menguji fungsionalitas seluruh antarmuka dan kesesuaian output",
    "Seluruh fitur CRUD, transaksi, dan laporan diuji secara menyeluruh",
    "Hasil: VALID \u2014 semua fitur berfungsi sesuai spesifikasi",
], font_size=9.5, color=CLR_GRAY)

# White-Box
add_shape_rect(slide, Inches(5.2), Inches(1.5), Inches(4.0), Pt(3), CLR_BLUE_DARK)
add_textbox(slide, Inches(5.2), Inches(1.65), Inches(4.0), Inches(0.35),
            "White-Box Testing (Basis Path)", font_size=14, bold=False, color=CLR_BLUE_DARK, font_name=FONT_TITLE)

add_bullet_items(slide, Inches(5.2), Inches(2.05), Inches(4.0), Inches(1.5), [
    "Fokus pada modul krusial: logika Penarikan Tabungan (fungsi tarikStore)",
    "Path 1: Saldo kurang \u2192 ditolak",
    "Path 2: Data tidak ditemukan \u2192 error handling",
    "Path 3: Transaksi valid \u2192 berhasil",
], font_size=9.5, color=CLR_GRAY)

# Bottom conclusion band
add_rounded_rect(slide, Inches(0.78), Inches(3.7), Inches(8.44), Inches(1.0), CLR_SNOW, CLR_CLOUD)
add_textbox(slide, Inches(1.0), Inches(3.8), Inches(8.0), Inches(0.3),
            "Kesimpulan Pengujian", font_size=13, bold=False, color=CLR_BLUE_MAIN, font_name=FONT_TITLE)
add_textbox(slide, Inches(1.0), Inches(4.1), Inches(8.0), Inches(0.5),
            "Seluruh skenario pengujian dinyatakan VALID dan BERHASIL. Algoritma penarikan tabungan dinyatakan aman dari cacat logika. Sistem siap digunakan secara operasional.",
            font_size=9.5, color=CLR_GRAY)

add_slide_number(slide, 10)


# ============================================================
# SLIDE 11: KESIMPULAN
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, CLR_WHITE)

add_page_title(slide, "Kesimpulan")

conclusions = [
    "Sistem berhasil mentransformasikan pencatatan manual menjadi digital terintegrasi, mempercepat pembuatan laporan, dan menekan angka human error.",
    "Fragmentasi data terselesaikan \u2014 manajemen SPP, tabungan, gaji, dan kas kini berada dalam satu wadah basis data.",
    "Transparansi dan akuntabilitas meningkat signifikan berkat Dashboard Analitik untuk yayasan dan riwayat pembayaran real-time bagi orang tua.",
    "Pengujian Black-Box dan White-Box membuktikan seluruh fitur berjalan valid dan algoritma aman dari cacat logika.",
]

colors = [CLR_BLUE_MAIN, CLR_BLUE_DARK, CLR_BLUE_MED, CLR_BLUE_LT]

for i, (conclusion, color) in enumerate(zip(conclusions, colors)):
    y = Inches(1.45) + Inches(i * 0.85)
    
    # Number badge
    add_number_badge(slide, Inches(0.78), y + Inches(0.05), i+1, color, Inches(0.35))
    
    # Conclusion text
    add_textbox(slide, Inches(1.3), y, Inches(7.9), Inches(0.75),
                conclusion, font_size=10.5, color=CLR_GRAY)

add_slide_number(slide, 11)


# ============================================================
# SLIDE 12: SARAN
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, CLR_WHITE)

add_page_title(slide, "Saran & Pengembangan ke Depan")

suggestions = [
    ("Payment Gateway", "Integrasi dengan Midtrans / Xendit untuk otomatisasi verifikasi pembayaran tanpa intervensi bendahara.", CLR_BLUE_MAIN),
    ("Aplikasi Mobile", "Pengembangan aplikasi Android/iOS agar orang tua dan siswa dapat memantau saldo dan tagihan dari smartphone.", CLR_BLUE_DARK),
    ("Notifikasi Otomatis", "Sistem pengingat jatuh tempo tagihan otomatis via WhatsApp atau Email kepada orang tua wali.", CLR_BLUE_MED),
]

for i, (title, desc, color) in enumerate(suggestions):
    y = Inches(1.45) + Inches(i * 1.1)
    
    # Accent bar on left
    add_shape_rect(slide, Inches(0.78), y, Pt(4), Inches(0.8), color)
    
    # Number
    add_textbox(slide, Inches(1.1), y, Inches(0.3), Inches(0.3),
                f"0{i+1}", font_size=18, bold=False, color=color, font_name=FONT_TITLE)
    
    # Title
    add_textbox(slide, Inches(1.6), y, Inches(7.5), Inches(0.3),
                title, font_size=13, bold=True, color=color, font_name=FONT_BODY)
    
    # Description
    add_textbox(slide, Inches(1.6), y + Inches(0.32), Inches(7.5), Inches(0.5),
                desc, font_size=10, color=CLR_GRAY)

add_slide_number(slide, 12)


# ============================================================
# SLIDE 13: TERIMA KASIH & Q&A
# ============================================================
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, CLR_CREAM)

# Decorative left bar (matching slide 1)
add_shape_rect(slide, Inches(0), Inches(0), Inches(0.15), Inches(5.63), CLR_BLUE_MAIN)

# Decorative circle (matching slide 1)
add_circle(slide, Inches(7.5), Inches(0.6), Inches(2.5), CLR_SNOW)
add_circle(slide, Inches(7.8), Inches(0.9), Inches(1.9), CLR_CLOUD)

# Thank you text
add_textbox(slide, Inches(0.78), Inches(1.4), Inches(6.5), Inches(0.8),
            "Terima Kasih", font_size=36, bold=False, color=CLR_BLACK,
            font_name=FONT_TITLE, alignment=PP_ALIGN.LEFT)

add_textbox(slide, Inches(0.78), Inches(2.15), Inches(5.0), Inches(0.35),
            "atas perhatian dan waktunya", font_size=13, color=CLR_BLUE_MAIN,
            font_name=FONT_BODY, italic=True)

# Divider
add_shape_rect(slide, Inches(0.78), Inches(2.7), Inches(2.5), Pt(1.5), CLR_CLOUD)

# Q&A
add_textbox(slide, Inches(0.78), Inches(2.9), Inches(5.0), Inches(0.4),
            "Sesi Tanya Jawab & Live Demo", font_size=16, bold=False,
            color=CLR_BLUE_DARK, font_name=FONT_TITLE)

add_textbox(slide, Inches(0.78), Inches(3.35), Inches(5.0), Inches(0.3),
            "Silakan ajukan pertanyaan", font_size=11, color=CLR_GRAY,
            font_name=FONT_BODY)

# Author info
add_shape_rect(slide, Inches(0.78), Inches(4.0), Inches(2.5), Pt(1), CLR_CLOUD)

add_multiline(slide, Inches(0.78), Inches(4.15), Inches(5.0), Inches(0.8), [
    ("Andri Gugun  |  NIM: [Isi NIM Anda]", 10, False, CLR_GRAY, FONT_BODY),
    ("Pondok Pesantren Latahzan Citeras", 9, False, CLR_LIGHT_GRAY, FONT_BODY),
], alignment=PP_ALIGN.LEFT, line_spacing=Pt(4))

add_slide_number(slide, 13)


# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           "Presentasi_Skripsi_Minimalist.pptx")
prs.save(output_path)
print(f"[OK] Presentasi berhasil dibuat: {output_path}")
print(f"     Total slide: {TOTAL_SLIDES}")
print(f"     Template: Minimalist Business Slides")
print(f"     Font: {FONT_TITLE} (judul) + {FONT_BODY} (body)")
