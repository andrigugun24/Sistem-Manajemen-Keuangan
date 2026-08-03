"""Inspect key slides of the Minimalist Business Slides template."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation(r"C:\wamp64\www\FinalProject\Minimalist Business Slides.pptx")

print(f"Slide: {prs.slide_width/914400:.2f}in x {prs.slide_height/914400:.2f}in")
print(f"Total slides: {len(prs.slides)}")
print(f"Total layouts: {len(prs.slide_layouts)}")
print()

# Only inspect first 15 slides to understand the design pattern
for slide_idx, slide in enumerate(prs.slides):
    if slide_idx >= 15:
        break
    
    layout_name = slide.slide_layout.name
    print(f"\n--- SLIDE {slide_idx + 1} (layout: '{layout_name}') ---")
    
    for shape in slide.shapes:
        shape_info = f"  {shape.name} | type={shape.shape_type}"
        shape_info += f" | pos=({shape.left/914400:.2f},{shape.top/914400:.2f})"
        shape_info += f" | size=({shape.width/914400:.2f},{shape.height/914400:.2f})"
        
        # Fill
        try:
            if hasattr(shape, "fill") and shape.fill.type is not None:
                try:
                    shape_info += f" | fill=#{shape.fill.fore_color.rgb}"
                except:
                    shape_info += f" | fill={shape.fill.type}"
        except:
            pass
        
        print(shape_info)
        
        # Text content
        if shape.has_text_frame:
            for p_idx, p in enumerate(shape.text_frame.paragraphs):
                if p.text.strip():
                    t = p.text[:80].encode('ascii','replace').decode()
                    align = str(p.alignment) if p.alignment else "inherit"
                    
                    font_info = ""
                    for run in p.runs:
                        f = run.font
                        parts = []
                        if f.size:
                            parts.append(f"{f.size/12700:.0f}pt")
                        if f.bold:
                            parts.append("BOLD")
                        if f.name:
                            parts.append(f.name)
                        try:
                            if f.color and f.color.type is not None:
                                parts.append(f"#{f.color.rgb}")
                        except:
                            parts.append("theme-color")
                        font_info = " | ".join(parts)
                        break
                    
                    print(f"    TEXT: '{t}' [{align}] ({font_info})")

# Also check key colors used in shapes across ALL slides
print("\n\n=== COLOR PALETTE (from shape fills across all slides) ===")
colors_found = set()
for slide in prs.slides:
    for shape in slide.shapes:
        try:
            if hasattr(shape, "fill") and shape.fill.type is not None:
                try:
                    c = str(shape.fill.fore_color.rgb)
                    colors_found.add(c)
                except:
                    pass
        except:
            pass

for c in sorted(colors_found):
    print(f"  #{c}")
